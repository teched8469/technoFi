from pptx import Presentation
from pptx.util import Inches

def get_slide_content(slide_number):
    print(f"Enter the title and content for slide {slide_number}:")
    title = input("Title: ")
    content = input("Content: ")
    return title, content

def get_image_path(slide_number):
    print(f"Enter the path to an image for slide {slide_number} (leave blank if no image):")
    image_path = input()
    return image_path if image_path.strip() else None

def create_presentation(num_slides, include_cover):
    # Create a new presentation object
    prs = Presentation()

    if include_cover:
        # Add a cover page
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = input("Enter the cover page title: ")
        subtitle.text = input("Enter the cover page subtitle: ")
    
    for i in range(1, num_slides + 1):
        # Prompt for slide content
        title, content = get_slide_content(i)
        
        # Choose a layout (1 is a title and content layout)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title and content
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        
        # Optionally, insert an image
        image_path = get_image_path(i)
        if image_path:
            left = top = Inches(1)
            slide.shapes.add_picture(image_path, left, top, height=Inches(4.5))
    
    # Save the presentation
    filename = input("Enter the filename: ")
    prs.save(filename)
    print(f"Presentation created successfully and saved as '{filename}'")

def main():
    try:
        num_slides = int(input("Enter the number of slides: "))
        if num_slides < 1:
            raise ValueError("Number of slides must be at least 1.")
        
        include_cover = input("Include a cover page? (yes/no): ").strip().lower() == 'yes'
        create_presentation(num_slides, include_cover)
    except ValueError as e:
        print(f"Invalid input: {e}")

if __name__ == "__main__":
    main()

